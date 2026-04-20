import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
from PIL import Image
import datetime

from pptx import Presentation
from pptx.util import Inches
from PIL import Image as PILImage
import os

def create_ppt(figures_with_titles, dept_groups_data=None):
    """
    Build a PPT matching the reference design:
      Slide 1 – 2x2 quadrant layout: blue title badge + black border + chart image
      Slide 2 – Large blue title + part-wise chart image
    The function re-draws all charts internally so they render cleanly at PPT size.
    """
    import io, tempfile, os
    import pandas as pd
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt
    import numpy as np
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    # ── Design constants ──────────────────────────────────────────────────────
    SLIDE_W  = Inches(10)
    SLIDE_H  = Inches(7.5)
    ACCENT   = RGBColor(0x4F, 0x81, 0xBD)
    WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
    BLACK    = RGBColor(0x00, 0x00, 0x00)
    BORDER_W = Pt(2.25)
    MARGIN   = Inches(0.45)
    GAP      = Inches(0.07)
    Q_W      = (SLIDE_W - 2 * MARGIN - GAP) / 2
    Q_H_TOP  = Inches(3.33)
    Q_H_BOT  = SLIDE_H - 2 * MARGIN - Q_H_TOP - GAP
    BADGE_H  = Inches(0.34)
    PAD      = Inches(0.06)

    # ── Helper: save matplotlib figure to BytesIO ─────────────────────────────
    def fig_to_buf(fig):
        buf = io.BytesIO()
        fig.savefig(buf, format='png', bbox_inches='tight', dpi=150)
        buf.seek(0)
        plt.close(fig)
        return buf

    # ── Helper: draw one quadrant (border + badge + image) ───────────────────
    def add_quadrant(slide, rl, rt, rw, rh, title_text, img_buf):
        # Border rectangle
        border = slide.shapes.add_shape(1, rl, rt, rw, rh)
        border.fill.background()
        border.line.color.rgb = BLACK
        border.line.width = BORDER_W
        # Title badge
        badge = slide.shapes.add_shape(1, rl+PAD, rt+PAD, rw-PAD*2, BADGE_H)
        badge.fill.solid()
        badge.fill.fore_color.rgb = ACCENT
        badge.line.fill.background()
        tf = badge.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = title_text
        run.font.size = Pt(18)
        run.font.color.rgb = WHITE
        # Chart image
        img_t = rt + PAD + BADGE_H + Inches(0.04)
        img_l = rl + PAD
        img_w = rw - PAD*2
        img_h = rh - (img_t - rt) - PAD
        slide.shapes.add_picture(img_buf, img_l, img_t, img_w, img_h)

    # ── Re-draw charts optimised for PPT quadrant size ────────────────────────
    # Pull the data from the passed figures' axes where possible,
    # but re-draw cleanly. We use fig_map for fallback to passed figs.
    fig_map = {t: f for t, f in figures_with_titles}

    # --- Chart 1: Production vs Rejection (render source figure directly) ---
    src1 = fig_map.get("Production vs Rejection")
    buf1 = fig_to_buf(src1) if src1 else None

    # --- Chart 2: Location-wise Rejection (render source figure directly) ---
    src2 = fig_map.get("Location-wise Rejection")
    buf2 = fig_to_buf(src2) if src2 else None

    # --- Chart 4: Customer-wise Rejection (render source figure directly) ---
    src4 = fig_map.get("Customer-wise Rejection")
    buf4 = fig_to_buf(src4) if src4 else None

    # --- Chart 5: Department-wise Rejection (horizontal bar, redrawn from raw data) ---
    dept_names_ppt  = ["Moulding", "Melting", "Technology", "Fettling", "Quality", "Maintenance"]
    dept_colors_ppt = {
        "Moulding": "#2E5B8A", "Melting": "#E09B3D",
        "Technology": "#3A8C5C", "Fettling": "#D4573C",
        "Quality": "#4A9FA8", "Maintenance": "#7B6FA3",
    }
    if dept_groups_data:
        fig5, axes5 = plt.subplots(3, 1, figsize=(5.0, 3.6))
        fig5.suptitle("Department-wise Rejection", fontsize=9, fontweight='bold', y=1.02)
        for ax5, (grp_label, prod_mt, dept_rej) in zip(axes5, dept_groups_data):
            total_rej_grp = sum(dept_rej.values())
            overall_pct   = (total_rej_grp / prod_mt * 100) if prod_mt > 0 else 0
            ax5.barh(0, prod_mt, height=0.5, color='#E8EDF2', zorder=2)
            left_val = 0.0
            for dept in dept_names_ppt:
                val = dept_rej.get(dept, 0)
                if val > 0:
                    ax5.barh(0, val, left=left_val, height=0.5,
                             color=dept_colors_ppt[dept], zorder=3)
                left_val += val
            offset = prod_mt * 0.015
            ax5.text(prod_mt + offset,  0.18, f"{prod_mt:.2f} MT",
                     va='center', fontsize=6.5, color='#444444')
            ax5.text(prod_mt + offset, -0.18,
                     f"Rej: {total_rej_grp:.3f} MT ({overall_pct:.2f}%)",
                     va='center', fontsize=6.5, color='#D4573C', fontweight='bold')
            breakdown = [f"{d}: {dept_rej.get(d,0):.3f}MT"
                         for d in dept_names_ppt if dept_rej.get(d, 0) > 0]
            if breakdown:
                ax5.text(0, -0.40, "  ".join(breakdown),
                         va='center', fontsize=5.5, color='#666666')
            ax5.set_yticks([0])
            ax5.set_yticklabels([grp_label], fontsize=7, fontweight='bold')
            ax5.set_ylim(-0.65, 0.65)
            ax5.set_xlim(0, prod_mt * 1.60)
            ax5.tick_params(axis='x', labelsize=6)
            ax5.spines[['top', 'right', 'left']].set_visible(False)
            ax5.yaxis.set_tick_params(length=0)
        from matplotlib.patches import Patch as MPatch
        legend_els = [MPatch(facecolor='#E8EDF2', label='Production')] + \
                     [MPatch(facecolor=dept_colors_ppt[d], label=d) for d in dept_names_ppt]
        fig5.legend(handles=legend_els, loc='lower center', ncol=4,
                    fontsize=5.5, bbox_to_anchor=(0.5, -0.06), frameon=False)
        fig5.tight_layout()
        buf5 = fig_to_buf(fig5)
    else:
        src5 = fig_map.get("Department-wise Rejection")
        buf5 = fig_to_buf(src5) if src5 else None

    # --- Chart 3: Part-wise (full slide) ---
    src3 = fig_map.get("Part-wise Production vs Rejection")
    if src3:
        fig_src = src3
        buf3 = fig_to_buf(fig_src)
    else:
        buf3 = None

    # ── Build presentation ────────────────────────────────────────────────────
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # ── SLIDE 1: 2x2 quadrant grid ────────────────────────────────────────────
    slide1 = prs.slides.add_slide(blank)
    LEFT_L  = MARGIN
    RIGHT_L = MARGIN + Q_W + GAP
    TOP_T   = MARGIN
    BOT_T   = MARGIN + Q_H_TOP + GAP

    quadrants = [
        (LEFT_L,  TOP_T, Q_W, Q_H_TOP, "Production vs Rejection",   buf1),
        (RIGHT_L, TOP_T, Q_W, Q_H_TOP, "Customer-wise Rejection",    buf4),
        (LEFT_L,  BOT_T, Q_W, Q_H_BOT, "Department-wise Rejection",  buf5),
        (RIGHT_L, BOT_T, Q_W, Q_H_BOT, "Location-wise Rejection",    buf2),
    ]
    for (ql, qt, qw, qh, qtitle, qbuf) in quadrants:
        if qbuf is not None:
            add_quadrant(slide1, ql, qt, qw, qh, qtitle, qbuf)

    # ── SLIDE 2: Part-wise full slide ─────────────────────────────────────────
    slide2 = prs.slides.add_slide(blank)
    # Title
    title_box = slide2.shapes.add_textbox(
        Inches(0.5), Inches(0.30), Inches(9.0), Inches(1.25))
    tf = title_box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "Part-wise Production vs Rejection"
    run.font.size = Pt(44)
    run.font.color.rgb = ACCENT
    # Chart image
    if buf3:
        slide2.shapes.add_picture(
            buf3,
            Inches(0.5), Inches(1.70),
            Inches(9.0), SLIDE_H - Inches(1.90)
        )

    # ── Save ──────────────────────────────────────────────────────────────────
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
    prs.save(tmp.name)
    return tmp.name
# ---------- UI ----------
st.set_page_config(layout="wide")
st.markdown('<style>div.block_container{padding-top:1rem;}</style>', unsafe_allow_html=True)

try:
    image = Image.open('logo.png')
    col1, col2 = st.columns([0.1, 0.9])
    with col1:
        st.image(image, width=100)
    with col2:
        st.title("Chandranchal BK - Production Dashboard")
except FileNotFoundError:
    st.title("Chandranchal BK - Production Dashboard")

col3, _, __ = st.columns([0.15, 0.45, 0.45])
with col3:
    st.write(f"Last updated: {datetime.datetime.now().strftime('%d %B %Y')}")

# ---------- File Upload ----------
file = st.file_uploader("Upload Production Excel File", type=["xlsx"])

if file:
    xl = pd.ExcelFile(file)

    # 🔥 Always pick FIRST sheet (most reliable for your files)
    sheet_name = xl.sheet_names[0]
    for s in xl.sheet_names:
        if any(month in s for month in [
            "January","February","March","April","May","June",
            "July","August","September","October","November","December",
            "Jan","Feb","Mar","Apr","Jun","Jul","Aug","Sep","Oct","Nov","Dec"
        ]):
            sheet_name = s
            break

    if sheet_name is None:
        sheet_name = xl.sheet_names[0]  # fallback
    # ── Load main data sheet (Row 4 = headers, data from row 5) ──────────────
    try:
        raw = pd.read_excel(file, sheet_name=sheet_name, header=None)

        header_row = None

        # 🔥 FIX: don't exceed available rows
        max_rows = min(15, len(raw))

        for i in range(max_rows):
            row = raw.iloc[i]

            for cell in row:
                cell_str = str(cell).lower().replace(" ", "").replace(".", "")

                if any(key in cell_str for key in ["sno", "slno", "srno"]):
                    header_row = i
                    break

            if header_row is not None:
                break

        # fallback
        if header_row is None:
            header_row = 3

        df = pd.read_excel(file, sheet_name=sheet_name, header=header_row)
    except Exception as e:
        st.error(f"❌ Could not read file's sheet: {e}")
        st.stop()

    # Rename duplicate rejection columns by location
    df = df.rename(columns={
        'Qty. Rej.':   'Foundry_Qty_Rej',
        'Wt in MT':    'Foundry_Wt_Rej',
        'Qty. Rej..1': 'MC_Qty_Rej',
        'Wt in MT.1':  'MC_Wt_Rej',
        'Qty. Rej..2': 'Cust_Qty_Rej',
        'Wt in MT.2':  'Cust_Wt_Rej',
        'Qty. Rej..3': 'Total_Qty_Rej',
        'Wt in MT.3':  'Total_Wt_Rej',
        # Department "Rej MT" columns (auto-named by pandas due to duplicates)
        'Rej MT':      'Moulding_Rej_MT',
        'Rej MT.1':    'Melting_Rej_MT',
        'Rej MT.2':    'Technology_Rej_MT',
        'Rej MT.3':    'Fettling_Rej_MT',
        'Rej MT.4':    'Quality_Rej_MT',
        'Rej MT.5':    'Maintenance_Rej_MT',
    })

    # ---------- Find S.No column dynamically ----------
    sno_col = None
    for col in df.columns:
        if 's.no' in str(col).lower().replace(" ", ""):
            sno_col = col
            break

    if sno_col is None:
        st.error("❌ Could not find S.No column. Check Excel format.")
        st.write("Detected columns:", df.columns)
        st.stop()

    # ---------- Filter valid rows ----------
    df = df[pd.to_numeric(df[sno_col], errors='coerce').notna()].copy()

    # Coerce numeric columns
    numeric_cols = ['Prod. Wt. in MT', 'Foundry_Wt_Rej', 'MC_Wt_Rej',
                    'Cust_Wt_Rej', 'Total_Wt_Rej',
                    'Moulding_Rej_MT', 'Melting_Rej_MT', 'Technology_Rej_MT',
                    'Fettling_Rej_MT', 'Quality_Rej_MT', 'Maintenance_Rej_MT']
    for c in numeric_cols:
        df[c] = pd.to_numeric(df[c], errors='coerce').fillna(0)

    st.success(f"✅ File uploaded — {len(df)} part records loaded")

    # ── KPI summary row ───────────────────────────────────────────────────────
    production = df['Prod. Wt. in MT'].sum()
    rejection  = df['Total_Wt_Rej'].sum()
    rej_pct    = (rejection / production * 100) if production > 0 else 0

    k1, k2, k3 = st.columns(3)
    k1.metric("Production (MT)", f"{production:,.3f}")
    k2.metric("Total Rejection (MT)", f"{rejection:,.3f}")
    k3.metric("Rejection %", f"{rej_pct:.2f}%")

    st.divider()

    # ── Chart layout ──────────────────────────────────────────────────────────
    left, right = st.columns(2)

    # ── 1. Production vs Rejection bar (group-wise + overall) ──────────────
    with left:
        st.subheader("📊 Production vs Rejection (MT)")

        # Read the bottom summary rows directly from raw Excel (header=None)
        # Row 423 = Pipe Fittings, Row 424 = Engineering (0-indexed after header=None)
        # ---------- READ SUMMARY DYNAMICALLY ----------
        raw_full = pd.read_excel(file, sheet_name=sheet_name, header=None)

        group_data = {}

        for i in range(raw_full.shape[0]):
            row = raw_full.iloc[i].astype(str)

            if row.str.contains("Pipe Fittings", case=False).any():
                g_prod = pd.to_numeric(raw_full.iloc[i, 2], errors='coerce')
                g_rej  = pd.to_numeric(raw_full.iloc[i, 3], errors='coerce')
                g_prod = float(g_prod) if pd.notna(g_prod) else 0.0
                g_rej  = float(g_rej)  if pd.notna(g_rej)  else 0.0
                group_data["Pipe Fittings"] = {
                    "prod": g_prod,
                    "rej":  g_rej,
                    "pct":  (g_rej / g_prod * 100) if g_prod > 0 else 0.0,
                }

            elif row.str.contains("Engineering", case=False).any():
                g_prod = pd.to_numeric(raw_full.iloc[i, 2], errors='coerce')
                g_rej  = pd.to_numeric(raw_full.iloc[i, 3], errors='coerce')
                g_prod = float(g_prod) if pd.notna(g_prod) else 0.0
                g_rej  = float(g_rej)  if pd.notna(g_rej)  else 0.0
                group_data["Engineering"] = {
                    "prod": g_prod,
                    "rej":  g_rej,
                    "pct":  (g_rej / g_prod * 100) if g_prod > 0 else 0.0,
                }

        # Add Overall row from df aggregates
        group_data["Overall"] = {
            "prod": production,
            "rej":  rejection,
            "pct":  rej_pct,
        }

        groups = ["Overall", "Pipe Fittings", "Engineering"]
        prods = [
            production,
            group_data.get("Pipe Fittings", {}).get("prod", 0),
            group_data.get("Engineering", {}).get("prod", 0),
        ]

        rejs = [
            rejection,
            group_data.get("Pipe Fittings", {}).get("rej", 0),
            group_data.get("Engineering", {}).get("rej", 0),
        ]

        pcts = [
            rej_pct,
            group_data.get("Pipe Fittings", {}).get("pct", 0),
            group_data.get("Engineering", {}).get("pct", 0),
        ]

        fig, ax = plt.subplots(figsize=(9, 5))
        y = range(len(groups))
        bar_h = 0.35

        prod_bars = ax.barh([i + bar_h/2 for i in y], prods, height=bar_h,
                            color='#2E5B8A', label='Production')
        rej_bars  = ax.barh([i - bar_h/2 for i in y], rejs,  height=bar_h,
                            color='#D4573C', label='Rejection')

        # Labels on production bars
        for i, (p, r, pct) in enumerate(zip(prods, rejs, pcts)):
            ax.text(p + max(prods)*0.01, i + bar_h/2,
                    f"{p:.2f} MT", va='center', fontsize=9, color='#2E5B8A')
            ax.text(r + max(prods)*0.01, i - bar_h/2,
                    f"{r:.2f} MT ({pct:.2f}%)", va='center', fontsize=9, color='#D4573C')

        ax.set_yticks(list(y))
        ax.set_yticklabels(groups, fontsize=10)
        ax.set_xlabel("Weight (MT)")
        ax.set_title("Production vs Rejection", fontsize=12)
        ax.legend(loc='lower right', fontsize=9)
        plt.tight_layout()

        st.pyplot(fig)
        chart1_fig = fig
        plt.close(fig)

    # ── 2. Location-wise Rejection pie ───────────────────────────────────────
    with right:
        st.subheader("📍 Location-wise Rejection")

        # --- Actual vhialues ---
        labels = ["Foundry", "M/c Shop", "Customer End"]

        weights = [
            df['Foundry_Wt_Rej'].sum(),
            df['MC_Wt_Rej'].sum(),
            df['Cust_Wt_Rej'].sum(),
        ]

        production = df['Prod. Wt. in MT'].sum()

        # --- Calculate ACTUAL percentages (NOT normalized) ---
        percentages = [
            (w / production * 100) if production > 0 else 0
            for w in weights
        ]

        # --- Remove zero entries ---
        filtered = [(l, w, p) for l, w, p in zip(labels, weights, percentages) if w > 0]

        if filtered:
            fl, fw, fp = zip(*filtered)

            fig, ax = plt.subplots()

            def make_autopct(actual_pcts, actual_wts):
                def my_autopct(pct):
                    idx = make_autopct.counter
                    make_autopct.counter += 1
                    return f"{actual_pcts[idx]:.2f}%\n({actual_wts[idx]:.2f} MT)"
                make_autopct.counter = 0
                return my_autopct

            PIE_COLORS = ['#2E5B8A', '#E09B3D', '#3A8C5C', '#D4573C', '#4A9FA8', '#7B6FA3']
            ax.pie(
                fw,
                labels=fl,
                autopct=make_autopct(fp, fw),
                startangle=140,
                colors=PIE_COLORS[:len(fw)],
                wedgeprops=dict(linewidth=0.8, edgecolor='white'),
                textprops=dict(fontsize=10)
            )

            # --- TOTAL DISPLAY ---
            total_wt = sum(fw)
            total_pct = sum(fp)

            plt.figtext(
                0.5, 0.02,
                f"Total Rejection = {total_pct:.2f}% ({total_wt:.2f} MT)",
                ha='center',
                fontsize=10
            )

            st.pyplot(fig)
            chart2_fig = fig
            plt.close(fig)

        else:
            st.info("No location rejection data available.")

    st.divider()

    left2, right2 = st.columns(2)

    # ── 3. Part-wise Rejection (top 10) ──────────────────────────────────────
    with left2:
        st.subheader("⚙️ Part-wise Production vs Rejection (All Parts)")
        part_group = (
            df.groupby('Part Name')[[
                'Prod. Wt. in MT',
                'Foundry_Wt_Rej',
                'MC_Wt_Rej',
                'Cust_Wt_Rej'
            ]]
            .sum()
            .sort_values('Prod. Wt. in MT', ascending=False)
        )

        # ✅ REMOVE ZERO PRODUCTION PARTS
        part_group = part_group[part_group['Prod. Wt. in MT'] > 0]

        part_group = part_group.fillna(0)
        part_group = part_group.iloc[::-1]

        fig_height = max(6, len(part_group) * 0.4)
        fig, ax = plt.subplots(figsize=(9, fig_height))

        y_pos = range(len(part_group))

        # ---------- BASE: Production ----------
        ax.barh(
            part_group.index,
            part_group['Prod. Wt. in MT'],
            color='#E8EDF2',
            label='Production'
        )

        # ---------- STACKED REJECTION ----------
        left_vals = [0] * len(part_group)

        colors = ['#2E5B8A', '#3A8C5C', '#D4573C']
        labels = ['Foundry', 'M/c Shop', 'Customer End']

        for col, color, label in zip(
            ['Foundry_Wt_Rej', 'MC_Wt_Rej', 'Cust_Wt_Rej'],
            colors,
            labels
        ):
            vals = part_group[col].values
            ax.barh(
                part_group.index,
                vals,
                left=left_vals,
                color=color,
                label=label
            )
            left_vals = [l + v for l, v in zip(left_vals, vals)]

        # ---------- LABELS ----------
        for i, (idx, row) in enumerate(part_group.iterrows()):
            prod = row['Prod. Wt. in MT']
            f = row['Foundry_Wt_Rej']
            m = row['MC_Wt_Rej']
            c = row['Cust_Wt_Rej']
            total = f + m + c

            # Avoid division issues
            rej_pct = (total / prod * 100) if prod > 0 else 0

            f_pct = (f / prod * 100) if prod > 0 else 0
            m_pct = (m / prod * 100) if prod > 0 else 0
            c_pct = (c / prod * 100) if prod > 0 else 0

            # Main label (production + rejection %)
            main_label = f"{prod:.1f} MT | Rej: {total:.1f} MT ({rej_pct:.2f}%)"

            # Breakdown
            sub_label = f"F:{f_pct:.2f}% M:{m_pct:.2f}% C:{c_pct:.2f}%"

            ax.text(prod + 0.1, i, main_label, va='center', fontsize=8)
            ax.text(prod + 0.1, i - 0.25, sub_label, va='center', fontsize=7, color='gray')

        ax.set_xlabel("Weight (MT)")
        ax.legend(loc='lower right', fontsize=8)
        ax.margins(y=0)
        plt.tight_layout()

        st.pyplot(fig)
        chart3_fig = fig
        plt.close(fig)

    # ── 4. Customer-wise Rejection pie ───────────────────────────────────────
    with right2:
        st.subheader("🏭 Customer-wise Rejection")

        cust_group = (
            df.groupby('Customer')['Total_Wt_Rej']
            .sum()
            .sort_values(ascending=False)
        )

        cust_group = cust_group[cust_group > 0]

        if not cust_group.empty:
            fig, ax = plt.subplots()

            # --- Values ---
            labels = cust_group.index.tolist()
            weights = cust_group.values.tolist()

            production = df['Prod. Wt. in MT'].sum()

            # --- ACTUAL % (not normalized) ---
            percentages = [
                (w / production * 100) if production > 0 else 0
                for w in weights
            ]

            def make_autopct(actual_pcts, actual_wts):
                def my_autopct(pct):
                    idx = make_autopct.counter
                    make_autopct.counter += 1
                    return f"{actual_pcts[idx]:.2f}%\n({actual_wts[idx]:.2f} MT)"
                make_autopct.counter = 0
                return my_autopct

            PIE_COLORS = ['#2E5B8A', '#E09B3D', '#3A8C5C', '#D4573C', '#4A9FA8', '#7B6FA3']
            wedges, texts = ax.pie(
                weights,
                labels=None,
                startangle=140,
                colors=PIE_COLORS[:len(weights)],
                wedgeprops=dict(linewidth=0.8, edgecolor='white')
            )

            import numpy as np

            for i, w in enumerate(wedges):
                angle = (w.theta2 + w.theta1) / 2
                x = np.cos(np.deg2rad(angle))
                y = np.sin(np.deg2rad(angle))

                ax.annotate(
                    f"{labels[i]}\n{percentages[i]:.2f}% ({weights[i]:.2f} MT)",
                    xy=(x, y),
                    xytext=(1.3 * x, 1.3 * y),
                    arrowprops=dict(arrowstyle="-", color='#aaaaaa', lw=0.8),
                    ha='center',
                    fontsize=9,
                    color='#333333'
                )

            # --- TOTAL DISPLAY ---
            total_wt = sum(weights)
            total_pct = sum(percentages)

            plt.figtext(
                0.5, 0.02,
                f"Total Rejection = {total_pct:.2f}% ({total_wt:.2f} MT)",
                ha='center',
                fontsize=10
            )

            ax.set_ylabel("")
            st.pyplot(fig)
            chart4_fig = fig
            plt.close(fig)

        else:
            st.info("No customer rejection data available.")

    st.divider()

    # ── 5. Department-wise Rejection – horizontal bar chart ──────────────────
    st.subheader("🏭 Department-wise Rejection (MT)")

    dept_names_list = ["Moulding", "Melting", "Technology", "Fettling", "Quality", "Maintenance"]
    # April layout: dept Rej MT values are at these raw summary-row col indices
    dept_raw_cols_apr = [7, 9, 11, 13, 15, 17]

    dept_colors = {
        "Moulding":    "#2E5B8A",
        "Melting":     "#E09B3D",
        "Technology":  "#3A8C5C",
        "Fettling":    "#D4573C",
        "Quality":     "#4A9FA8",
        "Maintenance": "#7B6FA3",
    }

    # ── Read group-level dept breakdowns from raw summary rows ────────────────
    raw_full_dept = pd.read_excel(file, sheet_name=sheet_name, header=None)

    def _get_group_dept(raw_df, group_keyword, raw_cols):
        for i in range(raw_df.shape[0]):
            row = raw_df.iloc[i].astype(str)
            if row.str.contains(group_keyword, case=False).any():
                vals = {}
                for name, col in zip(dept_names_list, raw_cols):
                    v = pd.to_numeric(raw_df.iloc[i, col], errors='coerce')
                    vals[name] = float(v) if pd.notna(v) else 0.0
                return vals
        return {name: 0.0 for name in dept_names_list}

    eng_dept = _get_group_dept(raw_full_dept, "Engineering",   dept_raw_cols_apr)
    pf_dept  = _get_group_dept(raw_full_dept, "Pipe Fittings", dept_raw_cols_apr)
    overall_dept = {name: eng_dept[name] + pf_dept[name] for name in dept_names_list}

    # Group production totals
    eng_prod        = group_data.get("Engineering",   {}).get("prod", 0)
    pf_prod         = group_data.get("Pipe Fittings", {}).get("prod", 0)
    total_prod_dept = df['Prod. Wt. in MT'].sum()

    groups_dept = [
        ("Overall",       total_prod_dept, overall_dept),
        ("Pipe Fittings", pf_prod,         pf_dept),
        ("Engineering",   eng_prod,        eng_dept),
    ]

    # ── Draw the chart ────────────────────────────────────────────────────────
    fig, axes = plt.subplots(3, 1, figsize=(11, 8), sharex=False)
    fig.suptitle("Department-wise Rejection by Group", fontsize=13, fontweight='bold', y=1.01)

    for ax, (grp_label, prod_mt, dept_rej) in zip(axes, groups_dept):
        total_rej_grp = sum(dept_rej.values())
        overall_pct   = (total_rej_grp / prod_mt * 100) if prod_mt > 0 else 0

        ax.barh(0, prod_mt, height=0.55, color='#E8EDF2', label='Production', zorder=2)

        left_val = 0.0
        for dept in dept_names_list:
            val = dept_rej.get(dept, 0)
            if val > 0:
                ax.barh(0, val, left=left_val, height=0.55,
                        color=dept_colors[dept], label=dept, zorder=3)
            left_val += val

        x_off = prod_mt * 0.012
        ax.text(prod_mt + x_off,  0.18,
                f"{prod_mt:.2f} MT", va='center', fontsize=9, color='#555555')
        ax.text(prod_mt + x_off, -0.18,
                f"Rej: {total_rej_grp:.3f} MT  ({overall_pct:.2f}%)",
                va='center', fontsize=9, color='#D4573C', fontweight='bold')

        breakdown_parts = []
        for dept in dept_names_list:
            val = dept_rej.get(dept, 0)
            pct = (val / prod_mt * 100) if prod_mt > 0 else 0
            if val > 0:
                breakdown_parts.append(f"{dept}: {val:.3f} MT ({pct:.2f}%)")
        if breakdown_parts:
            ax.text(0, -0.42, "  |  ".join(breakdown_parts),
                    va='center', fontsize=7.5, color='#555555')

        ax.set_yticks([0])
        ax.set_yticklabels([grp_label], fontsize=10, fontweight='bold')
        ax.set_ylim(-0.7, 0.7)
        ax.set_xlim(0, prod_mt * 1.55)
        ax.set_xlabel("Weight (MT)", fontsize=9)
        ax.tick_params(axis='x', labelsize=8)
        ax.spines[['top', 'right', 'left']].set_visible(False)
        ax.yaxis.set_tick_params(length=0)

    from matplotlib.patches import Patch
    legend_elements = [Patch(facecolor='#E8EDF2', label='Production')] + \
                      [Patch(facecolor=dept_colors[d], label=d) for d in dept_names_list]
    fig.legend(handles=legend_elements, loc='lower center', ncol=4,
               fontsize=8, bbox_to_anchor=(0.5, -0.04), frameon=False)

    plt.tight_layout()
    st.pyplot(fig)
    chart5_fig = fig
    plt.close(fig)


    st.divider()


    # ── 6. Material group summary (from Summary sheet) ────────────────────────
    st.subheader("🔩 Material Group Summary")
    try:
        raw = pd.read_excel(file, sheet_name="Summary", header=None)
        mat_data = {}
        for i in range(raw.shape[0]):
            cell = str(raw.iloc[i, 0]).strip()
            if cell in ("SG Iron", "Grey Iron"):
                prod_val    = pd.to_numeric(raw.iloc[i, 1], errors='coerce')
                rej_val     = pd.to_numeric(raw.iloc[i, 3], errors='coerce')
                rej_pct_val = pd.to_numeric(raw.iloc[i, 5], errors='coerce')
                mat_data[cell] = {
                    "Production (MT)": round(prod_val, 3) if not pd.isna(prod_val) else 0,
                    "Rejection (MT)":  round(rej_val, 3)  if not pd.isna(rej_val)  else 0,
                    "Rejection %":     f"{rej_pct_val*100:.2f}%" if not pd.isna(rej_pct_val) else "-",
                }
        if mat_data:
            st.table(pd.DataFrame(mat_data).T)
        else:
            st.info("Material group data not found in Summary sheet.")
    except Exception:
        st.info("Could not load material group summary.")

    # ── Raw data expander ─────────────────────────────────────────────────────
    with st.expander("📄 View Raw Data"):
        show_cols = ['S.No.', 'Customer', 'Part Name', 'Material grade',
                     'Prod. Wt. in MT', 'Foundry_Wt_Rej', 'MC_Wt_Rej',
                     'Cust_Wt_Rej', 'Total_Wt_Rej']
        available = [c for c in show_cols if c in df.columns]
        st.dataframe(df[available].reset_index(drop=True), use_container_width=True)

    st.divider()
    st.subheader("📥 Download Dashboard Report")

    if st.button("Generate PPT Report"):

        figures = [
            ("Production vs Rejection", chart1_fig),
            ("Location-wise Rejection", chart2_fig),
            ("Part-wise Production vs Rejection", chart3_fig),
            ("Customer-wise Rejection", chart4_fig),
            ("Department-wise Rejection", chart5_fig),
        ]

        ppt_file = create_ppt(figures, dept_groups_data=groups_dept)

        with open(ppt_file, "rb") as f:
            st.download_button(
                label="⬇️ Download PPT",
                data=f,
                file_name="Production_Dashboard_Report.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )

else:
    st.info("📂 Please upload a Production Excel file to get started.")
